# pip install PyPDF2
# pip install textract
# pip install bs4
# pip install python-pptx

import inspect
import logging
import pathlib
import re
from urllib.request import urlopen

import PyPDF2
import textract
from bs4 import BeautifulSoup
from pptx import Presentation

from constants import LOGGER_FORMAT, STATUS, MESSAGE, PDF_EXTENSION, TEXT_EXTENSION, PPTX_EXTENSION, HTML_EXTENSION


class Dotx2text:
    """ Class with single responsibility of extraction of text from a file """

    def __init__(self):
        """ Constructor """
        self.log = logging.getLogger(self.__class__.__name__)
        file_handler = logging.FileHandler(f'logs/{self.__class__.__name__}.log')
        formatter = logging.Formatter(LOGGER_FORMAT)
        file_handler.setFormatter(formatter)
        self.log.addHandler(file_handler)

    def get_text_from_file(self, filename):
        """
        Selects the correct decoder for the given file extension
        :param filename: String filename with extension
        :return: Clean Text
        """

        file_extension = pathlib.Path(filename).suffix
        directory = {PDF_EXTENSION: self.pdf2text, TEXT_EXTENSION: self.txt2text,
                     PPTX_EXTENSION: self.pptx2text, HTML_EXTENSION: self.html2text}
        if file_extension in directory:
            return directory[file_extension](filename)
        return self.generic2text(filename)

    def pdf2text(self, filename):
        """
        Extracts text from PDF files
        :param filename: Filename with .pdf extension
        :return: Clean Text
        """
        # Variable Declaration
        text = ""
        function_name = inspect.currentframe().f_code.co_name
        self.log.debug(f'{function_name} operation initiated.')

        # Text extraction from PDF
        try:
            pdf_file = open(filename, 'rb')
            pdf_reader = PyPDF2.PdfFileReader(pdf_file)
            for i in range(pdf_reader.numPages):
                page_object = pdf_reader.getPage(i)
                text += page_object.extractText()
        except Exception as e:
            return self.error_response(function_name, e)

        # Text cleaning
        text = text.replace("\r", "").replace("\n", "")
        self.log.debug(f'{function_name} operation successful.')
        return {STATUS: 100, MESSAGE: text}

    def generic2text(self, filename):
        """
        Extract text information from docx,csv,xls
        :param filename: Filename with correct extension
        :return: Clean Text
        """
        function_name = inspect.currentframe().f_code.co_name
        self.log.debug(f'{function_name} operation initiated.')

        # Text extraction from PDF
        try:
            byte_text = textract.process(filename)
            text = byte_text.decode("utf-8")
        except Exception as e:
            return self.error_response(function_name, e)

        # Text cleaning
        text = text.replace("\r", " ").replace("\n", " ").replace("\t", " ")
        self.log.debug(f'{function_name} operation successful.')
        return {STATUS: 100, MESSAGE: text}

    def txt2text(self, filename):
        """
        Extracts text from text files
        :param filename: Filename with .txt extension
        :return: Clean Text
        """
        function_name = inspect.currentframe().f_code.co_name
        self.log.debug(f'{function_name} operation initiated.')

        # Text extraction from Text file
        try:
            file = open(filename, encoding='utf-8')
            text = file.read().replace("\n", " ")
        except Exception as e:
            return self.error_response(function_name, e)

        self.log.debug(f'{function_name} operation successful.')
        return {STATUS: 100, MESSAGE: text}

    def pptx2text(self, filename):
        """
        Extracts text from Powerpoint files
        :param filename: Filename with .pptx extension
        :return: Clean Text
        """
        function_name = inspect.currentframe().f_code.co_name
        self.log.debug(f'{function_name} operation initiated.')
        # Variable Declaration
        result = ''

        # Text extraction from Powerpoint file
        try:
            prs = Presentation(filename)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        result += shape.text
        except Exception as e:
            return self.error_response(function_name, e)

        self.log.debug(f'{function_name} operation successful.')
        return {STATUS: 100, MESSAGE: result}

    def wiki2text(self, url):
        """
        Extracts text from wikipedia pages
        :param url: Web pages url
        :return: Clean Text
        """
        function_name = inspect.currentframe().f_code.co_name
        self.log.debug(f'{function_name} operation initiated.')

        # Text extraction from Wikipedia file
        try:
            source = urlopen(url).read()
            soup = BeautifulSoup(source, 'lxml')
        except Exception as e:
            return self.error_response(function_name, e)

        # Cleaning of Wikipedia page
        paras = []
        for paragraph in soup.find_all('p'):
            paras.append(str(paragraph.text))
        heads = []
        for head in soup.find_all('span', attrs={'mw-headline'}):
            heads.append(str(head.text))
        text = [val for pair in zip(paras, heads) for val in pair]
        text = ' '.join(text)
        text = re.sub(r"\[.*?\]+", '', text)
        text = text.replace('\n', '')[:-11]

        self.log.debug(f'{function_name} operation successful.')
        return {STATUS: 100, MESSAGE: text}

    def html2text(self, filename):
        """
        Extracts text from HTML files
        :param filename: Filename with .html extension
        :return: Clean Text
        """
        function_name = inspect.currentframe().f_code.co_name
        self.log.debug(f'{function_name} operation initiated.')
        # Text extraction from HTML file
        try:
            html = open(filename, encoding='utf-8')
            soup = BeautifulSoup(html, features="html.parser")
        except Exception as e:
            return self.error_response(function_name, e)

        # Cleaning of HTML page
        for script in soup(["script", "style"]):
            script.extract()
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip()
                  for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        self.log.debug(f'{function_name} operation successful')
        return {STATUS: 100, MESSAGE: text}

    def error_response(self, function_name, exception_message):
        """
        Generates and logs a error response in case of an file handling exception
        :param function_name: Name of the function
        :param exception_message: Exception message
        :return: error code for frontend
        """
        error_message = f'{function_name} Function failed. Error {exception_message}'
        self.log.error(error_message)
        return {STATUS: 400, MESSAGE: error_message}
