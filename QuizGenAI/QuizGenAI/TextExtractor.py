# pip install PyPDF2
# pip install textract
# pip install bs4
# pip install python-pptx

import PyPDF2
import textract
from pptx import Presentation
import glob
from urllib.request import urlopen
from bs4 import BeautifulSoup
import re


class Dotx2text:
    """ Sample Access Example
        from TextExtractor import Dotx2text 
        text = Dotx2text.functionname(arguments)"""

    @classmethod
    def dotpdf2text(cls, filename):
        """Extracts text from PDF files
           @input: (filename)-Filename with .pdf extension
           @Hyperparameters: ()-None
           @Output: (text)-Clean Text"""

        try:
            pdfFileObject = open(filename, 'rb')
            pdfReader = PyPDF2.PdfFileReader(pdfFileObject)
        except:
            return {"status": 400, "message": "dotpdf2text: Error in text extraction"}
        text = ""
        for i in range(pdfReader.numPages):
            try:
                pageObject = pdfReader.getPage(i)
                text += pageObject.extractText()
            except:
                return {"status": 400, "message": "dotpdf2text: Error in text extraction"}
        text = text.replace("\r", "")
        text = text.replace("\n", "")
        return {"status": 100, "message": text}

    @classmethod
    def generic2text(cls, filename):
        """Extract text information from docx,txt,pptx,csv,xls
           @input: (filename)-Filename with correct extension
           @Hyperparameters: ()-None
           @Output: (text)-Clean Text"""

        try:
            bytetext = textract.process(filename)
            text = bytetext.decode("utf-8")
        except:
            return {"status": 400, "message": "generic2text: Error in text extraction"}
        text = text.replace("\r", " ")
        text = text.replace("\n", " ")
        text = text.replace("\t", " ")
        return {"status": 100, "message": text}

    @classmethod
    def dottxt2text(cls, filename):
        """Extracts text from TXT files
           @input: (filename)-Filename with .txt extension
           @Hyperparameters: ()-None
           @Output: (text)-Clean Text"""

        try:
            file = open(filename, encoding='utf-8')
            text = file.read().replace("\n", " ")
        except:
            return {"status": 400, "message": "dottxt2text: Error in text extraction"}
        return {"status": 100, "message": text}

    @classmethod
    def dotpptx2text(cls, filename):
        """Extracts text from pptx files
           @input: (filename)-Filename with .pptx extension
           @Hyperparameters: ()-None
           @Output: (text)-Clean Text"""

        try:
            prs = Presentation(filename)
        except:
            return {"status": 400, "message": "dotpptx2text: Error in text extraction"}
        result = ''
        for slide in prs.slides:
            result += '\n 01newslide01 \n'
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    result += shape.text
        return {"status": 100, "message": result}

    @classmethod
    def wikiurl2text(cls, url):
        """Extracts text from WIKI pages files
           @input: (url)- Web pages
           @Hyperparameters: ()-None
           @Output: (text)-Clean Text"""

        try:
            source = urlopen(url).read()
            soup = BeautifulSoup(source, 'lxml')
        except:
            return {"status": 400, "message": "wikiurl2text: Error in text extraction"}
        paras = []
        for paragraph in soup.find_all('p'):
            paras.append(str(paragraph.text))
        heads = []
        for head in soup.find_all('span', attrs={'mw-headline'}):
            heads.append(str(head.text))
        text = [val for pair in zip(paras, heads) for val in pair]
        text = ' '.join(text)
        text = re.sub(r"\[.*?\]+", '', text)
        text = text.replace('\n', '')[:-11]
        return {"status": 100, "message": text}

    @classmethod
    def dothtml2text(cls, filename):
        """Extracts text from HTML files
           @input: (filename)- Filename with .html extension
           @Hyperparameters: ()-None
           @Output: (text)-Clean Text"""

        try:
            html = open(filename, encoding='utf-8')
            soup = BeautifulSoup(html, features="html.parser")
        except:
            return {"status": 400, "message": "dothtml2text: Error in text extraction"}
        for script in soup(["script", "style"]):
            script.extract()
        text = soup.get_text()

        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip()
                  for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        return {"status": 100, "message": text}


class DecoderSection:
    @classmethod
    def decoderselection(cls, filename):
        """Selects the correct decoder for the given file extension
           @input: (filename)-String filename with extension
           @HYperparamerter: ()-None
           @Output: (text)-Clean Text"""

        x = re.split("\.", filename)
        if x[-1].lower() == 'pdf':
            return Dotx2text.dotpdf2text(filename)
        elif x[-1].lower() == 'txt':
            return Dotx2text.dottxt2text(filename)
        elif x[-1].lower() == 'pptx':
            return Dotx2text.dotpptx2text(filename)
        elif x[-1].lower() == 'html':
            return Dotx2text.dothtml2text(filename)
        else:
            return Dotx2text.generic2text(filename)
